"""
PPTX (파워포인트/세미나 자료) 고품질 파서
슬라이드별 본문 + 표(Table) + 발표자 노트(Speaker Notes) 통합 추출
"""
from pathlib import Path
from typing import Dict, Any, List
from pptx import Presentation

class PptxParser:
    def parse(self, file_path: str | Path) -> Dict[str, Any]:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {path}")

        prs = Presentation(str(path))
        slides_data = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_content = []
            slide_tables = []
            speaker_notes = ""

            # 1. 도형 및 텍스트 상자 순회
            for shape in slide.shapes:
                # 텍스트가 있는 경우
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_content.append(f"- {text}")

                # 표(Table)가 있는 경우
                elif shape.has_table:
                    table = shape.table
                    md_table = self._parse_table_to_markdown(table)
                    if md_table:
                        slide_tables.append(md_table)

            # 2. 발표자 노트(Notes Slide) 추출
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    speaker_notes = notes_frame.text.strip()

            # 슬라이드별 Markdown 조합
            slide_md_parts = [f"## [Slide {idx}]"]
            if slide_content:
                slide_md_parts.append("\n".join(slide_content))
            if slide_tables:
                slide_md_parts.append("\n".join(slide_tables))
            if speaker_notes:
                slide_md_parts.append(f"\n> **[발표자 설명/노트]**\n> {speaker_notes}")

            slides_data.append({
                "slide_number": idx,
                "markdown": "\n\n".join(slide_md_parts)
            })

        full_markdown = "\n\n---\n\n".join([s["markdown"] for s in slides_data])

        return {
            "source_file": path.name,
            "file_type": "pptx",
            "total_slides": len(slides_data),
            "slides": slides_data,
            "markdown": full_markdown,
            "char_count": len(full_markdown)
        }

    def _parse_table_to_markdown(self, table) -> str:
        table_matrix = []
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                text = cell.text.strip().replace("\n", " ").replace("|", "\\|")
                row_cells.append(text if text else " ")
            table_matrix.append(row_cells)

        if not table_matrix:
            return ""

        max_cols = max(len(r) for r in table_matrix)
        normalized = [r + [" "] * (max_cols - len(r)) for r in table_matrix]

        header = normalized[0]
        divider = ["---"] * max_cols

        md_lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(divider) + " |"
        ]
        for r in normalized[1:]:
            md_lines.append("| " + " | ".join(r) + " |")

        return "\n".join(md_lines)
